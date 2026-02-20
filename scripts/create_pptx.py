"""
Script para criar apresentação PowerPoint do TVDEFleet
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Função para criar cor RGB
def RgbColor(r, g, b):
    """Retorna tupla RGB"""
    return (r, g, b)

def set_shape_color(shape, rgb_tuple):
    """Define cor de preenchimento de uma shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_to_pptx(rgb_tuple)

def rgb_to_pptx(rgb_tuple):
    """Converte tupla RGB para formato pptx"""
    from pptx.dml.color import RGBColor
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])

def set_font_color(font, rgb_tuple):
    """Define cor da fonte"""
    font.color.rgb = rgb_to_pptx(rgb_tuple)

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Cores do tema
AZUL_ESCURO = RgbColor(30, 58, 138)  # #1E3A8A
AZUL_CLARO = RgbColor(59, 130, 246)  # #3B82F6
VERDE = RgbColor(34, 197, 94)  # #22C55E
LARANJA = RgbColor(249, 115, 22)  # #F97316
CINZA = RgbColor(100, 116, 139)  # #64748B
BRANCO = RgbColor(255, 255, 255)
PRETO = RgbColor(15, 23, 42)  # #0F172A

def add_title_slide(prs, title, subtitle=""):
    """Adiciona slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fundo azul
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AZUL_ESCURO
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RgbColor(147, 197, 253)  # Azul claro
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, icon=""):
    """Adiciona slide com conteúdo em bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior azul
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_ESCURO
    bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = PRETO
        p.space_after = Pt(14)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, icon=""):
    """Adiciona slide com duas colunas"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_ESCURO
    bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    
    # Coluna esquerda
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = PRETO
        p.space_after = Pt(10)
    
    # Coluna direita
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = PRETO
        p.space_after = Pt(10)
    
    return slide

def add_table_slide(prs, title, headers, rows, icon=""):
    """Adiciona slide com tabela"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_ESCURO
    bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    
    # Tabela
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_width = Inches(12)
    table_height = Inches(0.5 * num_rows)
    left = Inches(0.666)
    top = Inches(1.8)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_CLARO
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(241, 245, 249)
    
    return slide


# ============== CRIAR SLIDES ==============

# Slide 1 - Capa
add_title_slide(prs, "TVDEFleet", "A Plataforma Completa de Gestão para Frotas TVDE")

# Slide 2 - O Problema
add_content_slide(prs, "Os Desafios da Gestão de Frotas TVDE", [
    "Dados dispersos em múltiplas plataformas (Uber, Bolt, etc.)",
    "Horas perdidas a inserir dados manualmente",
    "Documentação desorganizada de veículos e motoristas",
    "Falta de visibilidade sobre a performance real",
    "Dificuldade em escalar a operação",
    "Compliance complexo com múltiplas obrigações legais"
], "❌")

# Slide 3 - A Solução
add_content_slide(prs, "TVDEFleet: Tudo Num Só Lugar", [
    "Uma plataforma para toda a sua operação",
    "Importação automática de dados via RPA",
    "Documentação centralizada e sempre atualizada",
    "Dashboards em tempo real com KPIs relevantes",
    "Escalável de 5 a 500+ veículos",
    "Feito para Portugal com compliance integrado"
], "✅")

# Slide 4 - Gestão de Frota
add_table_slide(prs, "Gestão Completa de Frota", 
    ["Funcionalidade", "Benefício"],
    [
        ["Ficha Completa", "Todos os dados do veículo num só lugar"],
        ["Manutenção Preventiva", "Reduza avarias e custos inesperados"],
        ["Alertas de Seguros", "Nunca mais perca uma renovação"],
        ["Controlo de IPO", "Inspeções sempre em dia"],
        ["Gestão de Extintores", "Compliance de segurança garantido"],
        ["Histórico de Uso", "Saiba quem conduziu e quando"]
    ], "🚗")

# Slide 5 - Gestão de Motoristas
add_table_slide(prs, "Sistema de Classificações de Motoristas",
    ["Nível", "Requisitos", "Bónus"],
    [
        ["🥉 Bronze", "Início", "Base"],
        ["🥈 Prata", "3 meses + 60 pts", "+1%"],
        ["🥇 Ouro", "6 meses + 75 pts", "+2%"],
        ["💎 Platina", "12 meses + 85 pts", "+3.5%"],
        ["👑 Diamante", "24 meses + 95 pts", "+5%"]
    ], "👨‍💼")

# Slide 6 - Faturação
add_content_slide(prs, "Faturação Inteligente", [
    "Dashboard em Tempo Real - Receitas por plataforma (Uber, Bolt)",
    "Performance detalhada por motorista",
    "Relatórios Semanais Automáticos",
    "Exportação PDF profissional",
    "Cálculo automático de comissões",
    "Gestão de faturas de fornecedores"
], "💰")

# Slide 7 - RPA
add_table_slide(prs, "Automação RPA - Importação Automática",
    ["Plataforma", "Dados Importados"],
    [
        ["Uber", "Viagens, faturação, ganhos"],
        ["Bolt", "Corridas, valores, detalhes"],
        ["Prio", "Abastecimentos, litros, custos"],
        ["Via Verde", "Portagens, extratos mensais"]
    ], "🤖")

# Slide 8 - App Móvel
add_content_slide(prs, "Aplicação Móvel para Motoristas", [
    "Disponível na Google Play Store",
    "Dashboard de ganhos pessoal",
    "Envio rápido de recibos e despesas",
    "Notificações em tempo real",
    "Sistema de tickets/suporte",
    "Consulta de planos e documentação"
], "📱")

# Slide 9 - Planos
add_table_slide(prs, "Planos Flexíveis",
    ["Plano", "Ideal Para", "Inclui"],
    [
        ["Base Gratuito", "Começar", "Gestão básica"],
        ["Standard", "Frotas em crescimento", "+ Relatórios + Faturação"],
        ["Profissional", "Operações estabelecidas", "+ RPA + Módulos"],
        ["Enterprise", "Grandes frotas", "Tudo + Suporte VIP"]
    ], "📋")

# Slide 10 - Módulos
add_two_column_slide(prs, "Módulos Adicionais",
    [
        "✨ Autofaturação",
        "🔧 Manutenção Avançada", 
        "⚠️ Alertas de Custos",
        "📈 Dashboard de Ganhos"
    ],
    [
        "📊 Relatórios Detalhados",
        "🎯 Comissões Avançadas",
        "📅 Agenda Integrada",
        "💬 Mensagens Avançadas"
    ], "🧩")

# Slide 11 - Multi-utilizador
add_table_slide(prs, "Hierarquia de Utilizadores",
    ["Papel", "Acesso"],
    [
        ["Administrador", "Controlo total do sistema"],
        ["Gestor", "Gere múltiplos parceiros"],
        ["Parceiro", "Gestão da sua frota"],
        ["Contabilista", "Documentação financeira"],
        ["Inspetor", "Realiza vistorias"],
        ["Motorista", "App móvel e portal"]
    ], "👥")

# Slide 12 - Segurança
add_content_slide(prs, "Segurança e Compliance", [
    "Autenticação com tokens JWT seguros",
    "Permissões baseadas em papel",
    "Encriptação de dados em trânsito e repouso",
    "Backups automáticos diários",
    "RGPD compliant",
    "Logs de auditoria completos"
], "🔒")

# Slide 13 - Porquê TVDEFleet
add_two_column_slide(prs, "Porquê Escolher TVDEFleet?",
    [
        "✅ Tudo-em-Um",
        "    Uma plataforma, toda a gestão",
        "",
        "✅ Automação Real",
        "    RPA que poupa horas de trabalho",
        "",
        "✅ Feito para Portugal",
        "    Conhecemos o mercado TVDE"
    ],
    [
        "✅ Escalável",
        "    De 5 a 500+ veículos",
        "",
        "✅ Suporte Dedicado",
        "    Equipa sempre disponível",
        "",
        "✅ Atualizações Contínuas",
        "    Novas funcionalidades regularmente"
    ], "🎯")

# Slide 14 - Resultados
add_table_slide(prs, "Resultados Reais dos Nossos Clientes",
    ["Métrica", "Antes", "Depois", "Melhoria"],
    [
        ["Tempo em admin", "15h/semana", "3h/semana", "-80%"],
        ["Erros de dados", "Frequentes", "Raros", "-95%"],
        ["Docs em falta", "20%", "<2%", "-90%"],
        ["Visibilidade", "Limitada", "Total", "100%"]
    ], "📈")

# Slide 15 - Preços
add_table_slide(prs, "Investimento",
    ["Plano", "Mensal", "Inclui"],
    [
        ["Base", "Gratuito", "Gestão básica"],
        ["Standard", "€29.99", "+ Relatórios + Faturação"],
        ["Profissional", "€79.99", "+ RPA + Módulos"],
        ["Enterprise", "Sob consulta", "Tudo + Suporte VIP"]
    ], "💶")

# Slide 16 - Próximos Passos
add_content_slide(prs, "Comece Hoje - 3 Passos Simples", [
    "1️⃣ REGISTE-SE (5 minutos)",
    "2️⃣ CONFIGURE (30 minutos com apoio)",
    "3️⃣ COMECE A USAR (Imediato)",
    "",
    "✅ Trial gratuito de 30 dias",
    "✅ Onboarding assistido pela nossa equipa",
    "✅ Migração de dados incluída"
], "🚀")

# Slide 17 - Contactos
slide = add_title_slide(prs, "Fale Connosco", "")

# Adicionar info de contacto
contact_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(3))
tf = contact_box.text_frame
tf.word_wrap = True

contacts = [
    "📧 Email: geral@tvdefleet.com",
    "🌐 Website: www.tvdefleet.com", 
    "📱 App: Google Play Store",
    "",
    "Agende uma Demonstração Gratuita!"
]

for i, contact in enumerate(contacts):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = contact
    p.font.size = Pt(28) if i < 4 else Pt(32)
    p.font.color.rgb = BRANCO if i < 4 else RgbColor(74, 222, 128)
    p.font.bold = True if i >= 4 else False
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

# Guardar
output_path = "/app/docs/TVDEFleet_Apresentacao.pptx"
prs.save(output_path)
print(f"✅ PowerPoint criado com sucesso: {output_path}")
print(f"📊 Total de slides: {len(prs.slides)}")
